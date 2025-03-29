import os
import io
import tempfile
import re
import uuid
from PyPDF2 import PdfReader
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as PPTRGBColor
from pdfminer.high_level import extract_text, extract_pages
from pdfminer.layout import LAParams, LTTextBox, LTFigure, LTImage, LTRect, LTLine
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.pdfpage import PDFPage
from pdfminer.converter import TextConverter, PDFPageAggregator
from pdfminer.pdfparser import PDFParser
from pdfminer.pdfdocument import PDFDocument
from PIL import Image
import fitz  # PyMuPDF


def allowed_file(filename, allowed_extensions):
    """检查文件扩展名是否允许上传"""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in allowed_extensions


def clean_text_for_xml(text):
    """清理文本，移除XML不兼容的字符"""
    # 移除XML不兼容的控制字符
    text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]', '', text)
    # 替换其他可能导致问题的字符
    text = text.replace('\u0000', '')
    return text


def extract_images_from_pdf(pdf_path, output_dir=None):
    """从PDF提取图像"""
    if output_dir is None:
        output_dir = tempfile.mkdtemp()
    
    # 确保输出目录存在
    os.makedirs(output_dir, exist_ok=True)
    
    image_paths = []
    
    # 使用PyMuPDF提取图像
    try:
        doc = fitz.open(pdf_path)
        for page_index in range(len(doc)):
            page = doc[page_index]
            image_list = page.get_images(full=True)
            
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                image_ext = base_image["ext"]
                
                # 保存图像到临时文件
                image_filename = f"page{page_index+1}_img{img_index+1}.{image_ext}"
                image_path = os.path.join(output_dir, image_filename)
                
                with open(image_path, "wb") as img_file:
                    img_file.write(image_bytes)
                
                image_paths.append(image_path)
    except Exception as e:
        print(f"提取图像时出错: {e}")
    
    return image_paths


def get_pdf_layout_info(pdf_path):
    """获取PDF的布局信息，包括文本位置、图像、表格等"""
    layout_info = []
    
    try:
        # 打开PDF文件
        with open(pdf_path, 'rb') as f:
            # 初始化PDFMiner组件
            parser = PDFParser(f)
            document = PDFDocument(parser)
            rsrcmgr = PDFResourceManager()
            laparams = LAParams()
            device = PDFPageAggregator(rsrcmgr, laparams=laparams)
            interpreter = PDFPageInterpreter(rsrcmgr, device)
            
            # 处理每一页
            for page_num, page in enumerate(PDFPage.create_pages(document)):
                interpreter.process_page(page)
                layout = device.get_result()
                
                page_elements = []
                for element in layout:
                    if isinstance(element, LTTextBox):
                        # 文本元素
                        page_elements.append({
                            'type': 'text',
                            'x0': element.x0,
                            'y0': element.y0,
                            'x1': element.x1,
                            'y1': element.y1,
                            'text': element.get_text().strip(),
                            'font_size': get_font_size_estimate(element)
                        })
                    elif isinstance(element, LTFigure) or isinstance(element, LTImage):
                        # 图像或图形元素
                        page_elements.append({
                            'type': 'image',
                            'x0': element.x0,
                            'y0': element.y0,
                            'x1': element.x1,
                            'y1': element.y1
                        })
                    elif isinstance(element, LTRect) or isinstance(element, LTLine):
                        # 矩形或线条元素（可能是表格的一部分）
                        page_elements.append({
                            'type': 'rect',
                            'x0': element.x0,
                            'y0': element.y0,
                            'x1': element.x1,
                            'y1': element.y1
                        })
                
                # 按纵坐标排序，从上到下（注意PDF坐标系是从下到上的，所以要反转）
                page_elements.sort(key=lambda x: -x['y0'])
                
                layout_info.append(page_elements)
    except Exception as e:
        print(f"获取PDF布局信息时出错: {e}")
    
    return layout_info


def get_font_size_estimate(text_element):
    """估计文本元素的字体大小"""
    # 简单估计，可根据元素高度除以行数来计算
    try:
        height = text_element.y1 - text_element.y0
        lines = text_element.get_text().count('\n') + 1
        return round(height / lines)
    except:
        return 11  # 默认字体大小


def pdf_to_word(pdf_path, output_path=None):
    """将PDF转换为Word文档，保留格式和所有页面"""
    # 创建一个新的Word文档
    doc = Document()
    
    try:
        # 读取基本文本，作为备用
        fallback_text = clean_text_for_xml(extract_text(pdf_path))
        
        # 使用PyMuPDF处理所有页面
        pdf_doc = fitz.open(pdf_path)
        page_count = len(pdf_doc)
        
        for page_num in range(page_count):
            # 获取当前页面
            page = pdf_doc[page_num]
            
            # 提取页面文本
            page_text = page.get_text("text")
            page_text = clean_text_for_xml(page_text)
            
            # 只有在有多页的情况下才添加页码标题
            if page_count > 1 and page_num > 0:
                doc.add_page_break()
            
            # 处理页面文本
            if page_text.strip():
                paragraphs = page_text.split('\n\n')
                for para_text in paragraphs:
                    if para_text.strip():
                        para = doc.add_paragraph(para_text.strip())
            
            # 提取页面上的图像
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = pdf_doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # 保存到临时文件
                    temp_img = tempfile.NamedTemporaryFile(delete=False, suffix=f".{base_image['ext']}")
                    temp_img.write(image_bytes)
                    temp_img.close()
                    
                    # 添加到Word文档
                    doc.add_picture(temp_img.name, width=Inches(6))
                    
                    # 删除临时文件
                    try:
                        os.remove(temp_img.name)
                    except:
                        pass
                except Exception as e:
                    print(f"处理图像时出错: {e}")
        
        # 如果文档为空，使用备用文本
        if len(doc.paragraphs) == 0:
            paragraphs = fallback_text.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    doc.add_paragraph(para.strip())
    
    except Exception as e:
        # 添加一个错误信息段落
        doc.add_paragraph(f"PDF转换过程中遇到错误: {str(e)}")
        doc.add_paragraph("文档可能包含不兼容的内容，正在使用备用方法转换...")
        
        # 使用备用方法
        try:
            text = clean_text_for_xml(extract_text(pdf_path))
            paragraphs = text.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    doc.add_paragraph(para.strip())
        except:
            doc.add_paragraph("转换失败，请尝试其他格式。")
    
    # 如果没有指定输出路径，创建临时文件
    if output_path is None:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            output_path = tmp.name
    
    # 保存Word文档
    doc.save(output_path)
    
    return output_path


def pdf_to_excel(pdf_path, output_path=None):
    """将PDF转换为Excel表格，尝试识别表格结构"""
    # 创建一个新的Excel工作簿
    wb = Workbook()
    ws = wb.active
    ws.title = "Page 1"
    
    try:
        # 使用PyMuPDF直接访问PDF
        doc = fitz.open(pdf_path)
        page_count = len(doc)
        
        # 如果有多页，为每页创建一个工作表
        for page_num in range(1, page_count):
            wb.create_sheet(title=f"Page {page_num+1}")
        
        # 处理每一页
        for page_num in range(page_count):
            ws = wb.worksheets[page_num]
            page = doc[page_num]
            
            # 提取文本并按行分割
            page_text = page.get_text("text")
            page_text = clean_text_for_xml(page_text)
            lines = page_text.split('\n')
            
            # 按行填充Excel
            for i, line in enumerate(lines, start=1):
                if line.strip():
                    # 尝试根据分隔符分割单元格（制表符、多个空格等）
                    if '\t' in line:
                        # 如果有制表符，用它来分割
                        cells = line.split('\t')
                    else:
                        # 尝试智能分割行中的文本内容
                        cells = re.split(r'\s{2,}', line)  # 按两个或更多空格分割
                    
                    # 填充单元格
                    for j, cell_text in enumerate(cells, start=1):
                        cell_text = cell_text.strip()
                        if cell_text:
                            ws.cell(row=i, column=j, value=cell_text)
                            
                            # 添加格式
                            cell = ws.cell(row=i, column=j)
                            if i == 1:  # 第一行假设是标题
                                cell.font = Font(bold=True)
                                cell.alignment = Alignment(horizontal='center')
                                cell.fill = PatternFill(start_color="DDDDDD", end_color="DDDDDD", fill_type="solid")
            
            # 自动调整列宽
            for col in ws.columns:
                max_length = 0
                for cell in col:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                adjusted_width = max(max_length + 2, 10)
                ws.column_dimensions[col[0].column_letter].width = adjusted_width
    
    except Exception as e:
        # 使用备用方法
        ws = wb.active
        try:
            text = clean_text_for_xml(extract_text(pdf_path))
            lines = text.split('\n')
            for i, line in enumerate(lines, start=1):
                if line.strip():
                    ws.cell(row=i, column=1, value=line.strip())
        except:
            ws.cell(row=1, column=1, value="转换失败，请尝试其他格式。")
    
    # 如果没有指定输出路径，创建临时文件
    if output_path is None:
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
            output_path = tmp.name
    
    # 保存Excel文件
    wb.save(output_path)
    
    return output_path


def pdf_to_ppt(pdf_path, output_path=None):
    """将PDF转换为PowerPoint演示文稿，保留图像和格式"""
    # 创建一个新的演示文稿
    prs = Presentation()
    
    try:
        # 使用PyMuPDF直接访问PDF文档
        doc = fitz.open(pdf_path)
        page_count = len(doc)
        
        if page_count == 0:
            # 创建一个空白幻灯片
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf = txBox.text_frame
            tf.text = "PDF文档没有包含任何页面"
        else:
            # 处理每一页
            for page_num in range(page_count):
                page = doc[page_num]
                
                # 创建一个新的幻灯片
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)
                
                # 将页面渲染为图像
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                
                # 保存为临时图像文件
                temp_img = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                pix.save(temp_img.name)
                
                # 添加图像到幻灯片
                slide.shapes.add_picture(temp_img.name, 0, 0, width=prs.slide_width, height=prs.slide_height)
                
                # 提取页面文本作为注释
                page_text = page.get_text("text")
                page_text = clean_text_for_xml(page_text)
                
                # 添加文本框作为注释
                txBox = slide.shapes.add_textbox(Inches(0.1), Inches(6.5), Inches(9.8), Inches(0.5))
                tf = txBox.text_frame
                tf.text = f"第 {page_num+1} 页"
                
                # 添加隐藏的文本框用于搜索
                search_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(0.1), Inches(0.1))
                search_tf = search_box.text_frame
                search_tf.text = page_text
                search_p = search_tf.paragraphs[0]
                search_p.font.size = Pt(1)  # 非常小的字体
                search_p.font.color.rgb = PPTRGBColor(255, 255, 255)  # 白色，不可见
                
                # 删除临时文件
                try:
                    os.remove(temp_img.name)
                except:
                    pass
    
    except Exception as e:
        # 创建带有错误信息的幻灯片
        slide_layout = prs.slide_layouts[0]  # 标题幻灯片
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "PDF转换错误"
        
        subtitle = slide.placeholders[1]
        subtitle.text = f"转换过程中遇到错误: {str(e)}\n\n正在使用备用方法..."
        
        # 使用备用方法
        try:
            pdf = PdfReader(pdf_path)
            for i in range(len(pdf.pages)):
                page = pdf.pages[i]
                text = page.extract_text()
                text = clean_text_for_xml(text) if text else ""
                
                if text.strip():
                    slide_layout = prs.slide_layouts[1]  # 标题和内容
                    slide = prs.slides.add_slide(slide_layout)
                    
                    title = slide.shapes.title
                    title.text = f"第 {i+1} 页"
                    
                    content = slide.placeholders[1]
                    content.text = text
        except:
            # 如果备用方法也失败，添加一个失败信息幻灯片
            slide_layout = prs.slide_layouts[6]  # 空白布局
            slide = prs.slides.add_slide(slide_layout)
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            tf = txBox.text_frame
            tf.text = "无法提取PDF内容，请尝试其他格式。"
    
    # 如果没有指定输出路径，创建临时文件
    if output_path is None:
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            output_path = tmp.name
    
    # 保存演示文稿
    prs.save(output_path)
    
    return output_path


def pdf_to_markdown(pdf_path, output_path=None):
    """将PDF转换为Markdown文本，保留基本格式"""
    try:
        # 使用PyMuPDF直接访问PDF
        doc = fitz.open(pdf_path)
        page_count = len(doc)
        
        markdown_text = ""
        
        # 创建一个临时目录用于存储图像
        temp_img_dir = tempfile.mkdtemp()
        markdown_images_dir = "images"
        
        # 处理每一页
        for page_num in range(page_count):
            page = doc[page_num]
            
            # 添加页码标记
            markdown_text += f"# 第 {page_num + 1} 页\n\n"
            
            # 提取文本并进行清理
            page_text = page.get_text("text")
            page_text = clean_text_for_xml(page_text)
            
            # 按段落分割文本
            paragraphs = page_text.split('\n\n')
            
            # 处理每个段落
            for para in paragraphs:
                para = para.strip()
                if not para:
                    continue
                
                # 检测可能的标题
                lines = para.split('\n')
                for line in lines:
                    line = line.strip()
                    if not line:
                        continue
                    
                    # 尝试检测标题和列表项
                    if len(line) < 30 and not line.strip().endswith('.'):
                        if len(line) < 20:
                            markdown_text += f"## {line}\n\n"
                        else:
                            markdown_text += f"### {line}\n\n"
                    elif line.startswith('•') or line.startswith('-') or re.match(r'^\d+\.', line):
                        markdown_text += f"{line}\n"
                    else:
                        markdown_text += f"{line}\n\n"
            
            # 提取页面图像
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # 保存图像到临时文件
                    img_filename = f"page{page_num+1}_img{img_index+1}.{image_ext}"
                    img_path = os.path.join(temp_img_dir, img_filename)
                    
                    with open(img_path, "wb") as img_file:
                        img_file.write(image_bytes)
                    
                    # 添加图像引用到Markdown
                    markdown_text += f"![图片 {img_index+1}]({markdown_images_dir}/{img_filename})\n\n"
                except Exception as e:
                    print(f"处理图像时出错: {e}")
            
            # 添加页面分隔符
            markdown_text += "\n---\n\n"
    
    except Exception as e:
        # 使用备用方法
        markdown_text = f"转换过程中遇到错误: {str(e)}\n\n正在使用备用方法...\n\n"
        
        try:
            text = clean_text_for_xml(extract_text(pdf_path))
            lines = text.split('\n')
            current_line = ""
            
            for line in lines:
                line = line.strip()
                if not line:
                    if current_line:
                        markdown_text += current_line + "\n\n"
                        current_line = ""
                    continue
                
                # 检测可能的标题
                if len(line) < 50:
                    if line.endswith(':'):
                        markdown_text += f"## {line}\n\n"
                    elif len(line) < 30:
                        markdown_text += f"### {line}\n\n"
                    else:
                        current_line += line + " "
                else:
                    current_line += line + " "
            
            # 添加最后一行
            if current_line:
                markdown_text += current_line + "\n\n"
        except:
            markdown_text += "转换失败，请尝试其他格式。"
    
    # 如果没有指定输出路径，创建临时文件
    if output_path is None:
        with tempfile.NamedTemporaryFile(suffix='.md', delete=False) as tmp:
            output_path = tmp.name
    
    # 保存Markdown文件
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(markdown_text)
    
    # 清理临时文件
    try:
        import shutil
        shutil.rmtree(temp_img_dir)
    except:
        pass
    
    return output_path


def extract_summary(text, max_length=500):
    """从文本中提取摘要，返回指定长度内的摘要内容"""
    # 移除多余空白字符
    text = re.sub(r'\s+', ' ', text).strip()
    
    # 如果文本已经很短，直接返回
    if len(text) <= max_length:
        return text
    
    # 按句子分割文本
    sentences = re.split(r'(?<=[.!?。！？])\s+', text)
    
    # 初始化摘要和长度计数器
    summary = []
    current_length = 0
    
    # 逐句添加到摘要，直到达到最大长度
    for sentence in sentences:
        # 如果添加下一句会超出最大长度，则停止
        if current_length + len(sentence) + 1 > max_length:
            break
        
        summary.append(sentence)
        current_length += len(sentence) + 1  # +1 为空格
    
    # 如果一个完整的句子都无法添加，截取前max_length个字符
    if not summary:
        return text[:max_length] + "..."
    
    # 返回摘要
    result = " ".join(summary)
    
    # 如果摘要比原文短，添加省略号
    if len(result) < len(text):
        result += "..."
    
    return result
